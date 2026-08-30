#!/usr/bin/env python3
"""
Meerkat 文档工具链 v2: Markdown → HTML / PDF / Word / Excel / PPT
业界级排版 + 数据图表（真正的图表对象 / SVG 矢量，非静态图片）。

用法 (命令行):
  python3 doc_tools.py md2html <input.md> <output.html>
  python3 doc_tools.py md2docx <input.md> <output.docx>
  python3 doc_tools.py md2pdf  <input.md> <output.pdf>
  python3 doc_tools.py md2xlsx <input.md> <output.xlsx>
  python3 doc_tools.py md2pptx <input.md> <output.pptx>

用法 (函数调用, 供 LLM / bridge):
  from doc_tools import convert
  convert("html", markdown_text, "/path/out.html")   # 支持 html/docx/pdf/xlsx/pptx

图表 DSL (在 Markdown 中内嵌 fenced code block, 语言标签为 chart):
  ```chart
  {"type": "bar", "title": "月度销量", "categories": ["1月","2月","3月"],
   "series": [{"name": "销量", "data": [120, 150, 180]}]}
  ```

支持的图表 type:
  bar / barh / line / area / pie / donut / radar / scatter / funnel / gauge
"""

import base64
import io
import json
import os
import re
import sys
import tempfile

# ======================================================================
# 常量: 字体 / 配色
# ======================================================================

# 字体目录: fonts/ = 从 Noto ttc 提取的 SC 子字体 (CFF 轮廓, matplotlib 支持)
#           fonts_ttf/ = cu2qu 转换后的 TrueType 轮廓 (reportlab 支持, 由 convert_fonts.py 生成)
_FONT_DIR = "/home/chinux/jupyterlab/meerkatai/fonts"
_FONT_DIR_TTF = "/home/chinux/jupyterlab/meerkatai/fonts_ttf"


def _resolve_mpl_fonts():
    """matplotlib 用字体: 返回 (regular, bold, family)。支持 CFF 轮廓。"""
    candidates = [
        (os.path.join(_FONT_DIR, "NotoSansCJKsc-Regular.ttf"),
         os.path.join(_FONT_DIR, "NotoSansCJKsc-Bold.ttf"), "Noto Sans CJK SC"),
        ("/System/Library/Fonts/PingFang.ttc", "/System/Library/Fonts/PingFang.ttc", "PingFang SC"),
        ("C:/Windows/Fonts/msyh.ttc", "C:/Windows/Fonts/msyhbd.ttc", "Microsoft YaHei"),
    ]
    for reg, bold, fam in candidates:
        if reg and os.path.exists(reg):
            return reg, bold, fam
    return None, None, "Helvetica"


def _resolve_pdf_fonts():
    """reportlab 用字体: 返回 (regular, bold, serif_regular, serif_bold)。必须 TrueType 轮廓。"""
    candidates = [
        # 1. cu2qu 转换后的 Noto TrueType (业界顶级, 由 convert_fonts.py 生成)
        (os.path.join(_FONT_DIR_TTF, "NotoSansCJKsc-Regular.ttf"),
         os.path.join(_FONT_DIR_TTF, "NotoSansCJKsc-Bold.ttf"),
         os.path.join(_FONT_DIR_TTF, "NotoSerifCJKsc-Regular.ttf"),
         os.path.join(_FONT_DIR_TTF, "NotoSerifCJKsc-Bold.ttf")),
        # 2. 系统 TrueType 中文字体 (立即可用): Droid 黑体 + 文鼎明体
        ("/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
         "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
         "/usr/share/fonts/truetype/arphic/uming.ttc",
         "/usr/share/fonts/truetype/arphic/uming.ttc"),
    ]
    for reg, bold, sreg, sbold in candidates:
        if reg and os.path.exists(reg):
            return reg, bold, sreg, sbold
    return None, None, None, None


# 专业数据可视化配色 (AntV G2 经典色板)
PALETTE = [
    "#5B8FF9", "#61DDAA", "#65789B", "#F6BD16", "#7262FD",
    "#78D3F8", "#9661BC", "#F6903D", "#008685", "#F08BB4",
]

# HTML 主题色
_HTML_CSS = """
:root {
  --primary: #2563eb;
  --primary-soft: #eff6ff;
  --heading: #1e293b;
  --text: #334155;
  --muted: #64748b;
  --border: #e2e8f0;
  --bg-soft: #f8fafc;
  --accent: #0ea5e9;
}
* { box-sizing: border-box; }
body {
  margin: 0;
  font-family: "Noto Sans CJK SC", "PingFang SC", "Microsoft YaHei", "Segoe UI", system-ui, sans-serif;
  color: var(--text);
  line-height: 1.8;
  background: #fff;
  font-size: 15px;
  -webkit-font-smoothing: antialiased;
}
/* ===== 封面 ===== */
.cover {
  padding: 96px 48px 72px;
  background: linear-gradient(135deg, #1e3a8a 0%, #2563eb 55%, #0ea5e9 100%);
  color: #fff;
  text-align: center;
  page-break-after: always;
}
.cover h1 { font-size: 36px; margin: 0 0 16px; font-weight: 700; letter-spacing: 1px; }
.cover .subtitle { font-size: 16px; opacity: .9; margin: 0 0 40px; }
.cover .meta { display: inline-block; padding: 6px 18px; border: 1px solid rgba(255,255,255,.5); border-radius: 999px; font-size: 13px; opacity: .95; }
/* ===== 目录 ===== */
.toc {
  max-width: 860px; margin: 0 auto; padding: 40px 32px;
}
.toc h2 { font-size: 20px; color: var(--heading); border-bottom: 2px solid var(--primary); padding-bottom: 8px; }
.toc ul { list-style: none; padding: 0; }
.toc li { margin: 6px 0; }
.toc a { color: var(--text); text-decoration: none; display: block; padding: 8px 14px; border-left: 3px solid transparent; border-radius: 0 6px 6px 0; transition: all .15s; }
.toc a:hover { background: var(--primary-soft); border-left-color: var(--primary); color: var(--primary); }
.toc .lvl2 { padding-left: 34px; font-size: 14px; color: var(--muted); }
/* ===== 正文 ===== */
main { max-width: 860px; margin: 0 auto; padding: 16px 32px 80px; }
h1, h2, h3, h4, h5, h6 { color: var(--heading); line-height: 1.4; margin: 1.6em 0 .6em; }
h1 { font-size: 28px; border-bottom: 3px solid var(--primary); padding-bottom: 10px; }
h2 { font-size: 22px; border-bottom: 1px solid var(--border); padding-bottom: 8px; }
h3 { font-size: 18px; }
h4 { font-size: 16px; }
p { margin: .7em 0; }
strong { color: var(--heading); }
a { color: var(--primary); }
code {
  background: var(--bg-soft); border: 1px solid var(--border); border-radius: 4px;
  padding: 1px 6px; font-family: "SFMono-Regular", Consolas, "Liberation Mono", monospace; font-size: 13px; color: #be123c;
}
pre {
  background: #0f172a; color: #e2e8f0; padding: 18px 20px; border-radius: 8px;
  overflow-x: auto; font-size: 13px; line-height: 1.6;
}
pre code { background: none; border: none; color: inherit; padding: 0; }
blockquote {
  margin: 1em 0; padding: 10px 20px; border-left: 4px solid var(--accent);
  background: var(--bg-soft); color: var(--muted); border-radius: 0 6px 6px 0;
}
table { border-collapse: collapse; width: 100%; margin: 1.2em 0; font-size: 14px; }
th, td { border: 1px solid var(--border); padding: 9px 12px; text-align: left; }
th { background: var(--primary); color: #fff; font-weight: 600; }
tr:nth-child(even) td { background: var(--bg-soft); }
tr:hover td { background: var(--primary-soft); }
img { max-width: 100%; border-radius: 6px; }
hr { border: none; border-top: 1px solid var(--border); margin: 2em 0; }
/* ===== 图表 ===== */
.figure {
  margin: 1.6em 0; padding: 16px; border: 1px solid var(--border); border-radius: 8px;
  background: #fff; text-align: center; box-shadow: 0 1px 3px rgba(0,0,0,.05);
}
.figure svg { max-width: 100%; height: auto; }
/* ===== 打印 ===== */
@media print {
  .cover { page-break-after: always; }
  .toc { page-break-after: always; }
  body { font-size: 12pt; }
  main { max-width: 100%; }
  a { color: inherit; }
  .figure { box-shadow: none; page-break-inside: avoid; }
}
@page { margin: 20mm; }
"""


# ======================================================================
# Markdown 解析器 (块级 + 行内)
# ======================================================================

def _strip_inline(text):
    """去除行内标记, 得到纯文本。"""
    text = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text


def _inline_to_html(text):
    """行内标记 → HTML (转义 + 加粗/斜体/代码/链接)。"""
    def esc(s):
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    text = esc(text)
    # 先处理图片和链接
    text = re.sub(r"!\[([^\]]*)\]\(([^)]*)\)", r'<img src="\2" alt="\1">', text)
    text = re.sub(r"\[([^\]]*)\]\(([^)]*)\)", r'<a href="\2">\1</a>', text)
    # 加粗 / 斜体 / 行内代码
    text = re.sub(r"\*\*([^*]+)\*\*", r"<strong>\1</strong>", text)
    text = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"<em>\1</em>", text)
    text = re.sub(r"`([^`]+)`", r"<code>\1</code>", text)
    return text


def _parse_blocks(md_text):
    """把 markdown 解析成结构化块: (type, payload)。"""
    blocks = []
    lines = md_text.split("\n")
    i, n = 0, len(lines)

    def _chart_or_code(line, i):
        # fenced code block: ```chart ... ``` 或 ```lang ... ```
        m = re.match(r"^```\s*(\w*)\s*$", line.strip())
        if not m:
            return None
        lang = m.group(1)
        buf = []
        i += 1
        while i < n and not lines[i].strip().startswith("```"):
            buf.append(lines[i])
            i += 1
        if i < n:
            i += 1  # 跳过结尾 ```
        content = "\n".join(buf)
        if lang == "chart":
            try:
                return ("chart", json.loads(content)), i
            except Exception:
                return ("code", content), i
        return ("code", content), i

    while i < n:
        line = lines[i].rstrip()
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # fenced code / chart
        if stripped.startswith("```"):
            res = _chart_or_code(line, i)
            if res:
                blocks.append(res[0])
                i = res[1]
                continue

        # 标题
        m = re.match(r"^(#{1,6})\s+(.*)", stripped)
        if m:
            blocks.append(("heading", (len(m.group(1)), m.group(2).strip())))
            i += 1
            continue

        # 图片 (Markdown 内联 base64 图)
        m = re.match(r"!\[([^\]]*)\]\(data:image/([a-zA-Z]+);base64,([A-Za-z0-9+/=]+)\)", stripped)
        if m:
            blocks.append(("image", (m.group(1), m.group(2), m.group(3))))
            i += 1
            continue

        # 表格
        if stripped.startswith("|"):
            rows = []
            while i < n and lines[i].strip().startswith("|"):
                rows.append(lines[i].strip())
                i += 1
            blocks.append(("table", rows))
            continue

        # 列表 (含嵌套, 只支持一级缩进)
        if re.match(r"^\s*([-*+]|\d+\.)\s+", stripped):
            items = []
            while i < n and re.match(r"^\s*([-*+]|\d+\.)\s+", lines[i].strip()):
                items.append(re.sub(r"^\s*([-*+]|\d+\.)\s+", "", lines[i].strip()))
                i += 1
            blocks.append(("list", items))
            continue

        # 引用
        if stripped.startswith(">"):
            quotes = []
            while i < n and lines[i].strip().startswith(">"):
                quotes.append(lines[i].strip().lstrip(">").strip())
                i += 1
            blocks.append(("quote", " ".join(quotes)))
            continue

        # 分隔线
        if re.match(r"^(-{3,}|\*{3,}|_{3,})$", stripped):
            blocks.append(("hr", None))
            i += 1
            continue

        # 段落
        para = []
        while i < n and lines[i].strip() and not re.match(
                r"^(#{1,6}\s|\||\s*[-*+]\s|\s*\d+\.\s|!\[|>|```|-{3,}|\*{3,}|_{3,})", lines[i]):
            para.append(lines[i].strip())
            i += 1
        blocks.append(("para", " ".join(para)))

    return blocks


def _parse_table_rows(rows):
    """解析 markdown 表格行为 header + data 行 (跳过 |---| 分隔行)。"""
    header, data = None, []
    for r in rows:
        cells = [c.strip() for c in r.strip().strip("|").split("|")]
        if all(re.match(r"^:?-{2,}:?$", c) for c in cells):
            continue
        if header is None:
            header = cells
        else:
            data.append(cells)
    return header, data


def _extract_title(blocks):
    """取第一个 H1 作为文档标题。"""
    for typ, payload in blocks:
        if typ == "heading" and payload[0] == 1:
            return _strip_inline(payload[1])
    return "文档"


# ======================================================================
# 图表引擎 (matplotlib)
# ======================================================================

_mpl_ready = False


def _mpl_setup():
    """惰性初始化 matplotlib (字体 + 样式), 只做一次。"""
    global _mpl_ready
    if _mpl_ready:
        return True
    try:
        import matplotlib
        matplotlib.use("Agg")
        from matplotlib import font_manager
        import matplotlib.pyplot as plt

        reg, bold, fam = _resolve_mpl_fonts()
        registered = set()
        for p in (reg, bold):
            if p and os.path.exists(p) and p not in registered:
                try:
                    font_manager.fontManager.addfont(p)
                    registered.add(p)
                except Exception:
                    pass

        plt.rcParams.update({
            "font.sans-serif": [fam, "DejaVu Sans"],
            "font.family": "sans-serif",
            "axes.unicode_minus": False,
            "figure.dpi": 110,
            "svg.fonttype": "path",
            "axes.edgecolor": "#cbd5e1",
            "axes.linewidth": 0.8,
            "axes.grid": True,
            "grid.color": "#e2e8f0",
            "grid.linewidth": 0.6,
            "axes.axisbelow": True,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.titlesize": 14,
            "axes.titleweight": "bold",
            "axes.labelsize": 11,
            "xtick.labelsize": 10,
            "ytick.labelsize": 10,
            "legend.frameon": False,
            "legend.fontsize": 10,
        })
        _mpl_ready = True
        return True
    except Exception:
        return False


def _chart_axes(chart, fig):
    """根据图表高度返回 axes。"""
    h = float(chart.get("height") or 0)
    if h <= 0:
        # 默认高度
        n = len(chart.get("categories") or []) or 1
        h = max(3.2, min(8.0, 3.0 + n * 0.35))
    fig.set_size_inches(7.2, h)
    return fig.subplots()


def _bar_chart(ax, chart):
    import numpy as np
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    x = np.arange(len(cats))
    width = 0.7 / max(len(series), 1)
    unit = chart.get("unit") or ""
    for si, s in enumerate(series):
        data = [float(v) for v in s.get("data", [])]
        off = (si - (len(series) - 1) / 2) * width
        bars = ax.bar(x + off, data, width=width, label=s.get("name", ""),
                      color=PALETTE[si % len(PALETTE)], zorder=3)
        for b, v in zip(bars, data):
            ax.text(b.get_x() + b.get_width() / 2, b.get_height(),
                    f"{v:g}{unit}", ha="center", va="bottom", fontsize=8)
    ax.set_xticks(x)
    ax.set_xticklabels(cats)
    if chart.get("x_label"):
        ax.set_xlabel(chart["x_label"])
    if chart.get("y_label"):
        ax.set_ylabel(chart["y_label"])
    if series and series[0].get("name"):
        ax.legend()


def _barh_chart(ax, chart):
    import numpy as np
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    y = np.arange(len(cats))
    height = 0.7 / max(len(series), 1)
    unit = chart.get("unit") or ""
    for si, s in enumerate(series):
        data = [float(v) for v in s.get("data", [])]
        off = (si - (len(series) - 1) / 2) * height
        bars = ax.barh(y + off, data, height=height, label=s.get("name", ""),
                       color=PALETTE[si % len(PALETTE)], zorder=3)
        for b, v in zip(bars, data):
            ax.text(b.get_width(), b.get_y() + b.get_height() / 2,
                    f" {v:g}{unit}", va="center", ha="left", fontsize=8)
    ax.set_yticks(y)
    ax.set_yticklabels(cats)
    if chart.get("x_label"):
        ax.set_xlabel(chart["x_label"])
    if series and series[0].get("name"):
        ax.legend()


def _line_chart(ax, chart):
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    x = list(range(len(cats)))
    unit = chart.get("unit") or ""
    for si, s in enumerate(series):
        data = [float(v) for v in s.get("data", [])]
        ax.plot(x, data, marker="o", linewidth=2, markersize=5,
                label=s.get("name", ""), color=PALETTE[si % len(PALETTE)], zorder=3)
        for xi, v in zip(x, data):
            ax.text(xi, v, f"{v:g}{unit}", ha="center", va="bottom", fontsize=8)
    ax.set_xticks(x)
    ax.set_xticklabels(cats)
    if chart.get("x_label"):
        ax.set_xlabel(chart["x_label"])
    if chart.get("y_label"):
        ax.set_ylabel(chart["y_label"])
    if series and series[0].get("name"):
        ax.legend()


def _area_chart(ax, chart):
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    x = list(range(len(cats)))
    unit = chart.get("unit") or ""
    for si, s in enumerate(series):
        data = [float(v) for v in s.get("data", [])]
        ax.fill_between(x, data, alpha=0.25, color=PALETTE[si % len(PALETTE)])
        ax.plot(x, data, linewidth=2, label=s.get("name", ""),
                color=PALETTE[si % len(PALETTE)], zorder=3)
    ax.set_xticks(x)
    ax.set_xticklabels(cats)
    if chart.get("x_label"):
        ax.set_xlabel(chart["x_label"])
    if chart.get("y_label"):
        ax.set_ylabel(chart["y_label"])
    if series and series[0].get("name"):
        ax.legend()


def _pie_chart(ax, chart, donut=False):
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    values = [float(v) for v in (series[0].get("data", []) if series else [])]
    if not values:
        return
    total = sum(values) or 1
    colors = [PALETTE[i % len(PALETTE)] for i in range(len(values))]
    if donut:
        wedges, _texts, _autotexts = ax.pie(values, labels=cats, colors=colors, startangle=90,
                                            wedgeprops=dict(width=0.42, edgecolor="white", linewidth=2),
                                            autopct=lambda p: f"{p:.1f}%")
    else:
        wedges, _texts, _autotexts = ax.pie(values, labels=cats, colors=colors, startangle=90,
                                            autopct=lambda p: f"{p:.1f}%")
    for w in wedges:
        w.set_edgecolor("white")
        w.set_linewidth(1.5)
    ax.axis("equal")


def _radar_chart(ax, chart):
    import numpy as np
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    N = len(cats)
    if N < 3:
        return
    angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
    angles += angles[:1]
    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)
    for si, s in enumerate(series):
        data = [float(v) for v in s.get("data", [])]
        if len(data) != N:
            continue
        vals = data + data[:1]
        ax.plot(angles, vals, linewidth=2, label=s.get("name", ""),
                color=PALETTE[si % len(PALETTE)])
        ax.fill(angles, vals, alpha=0.2, color=PALETTE[si % len(PALETTE)])
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(cats, fontsize=10)
    if series and series[0].get("name"):
        ax.legend(loc="upper right", bbox_to_anchor=(1.25, 1.1))


def _scatter_chart(ax, chart):
    series = chart.get("series") or []
    for si, s in enumerate(series):
        pts = s.get("data", [])
        xs = [float(p[0]) for p in pts if isinstance(p, (list, tuple)) and len(p) >= 2]
        ys = [float(p[1]) for p in pts if isinstance(p, (list, tuple)) and len(p) >= 2]
        if xs:
            ax.scatter(xs, ys, s=55, alpha=0.8, label=s.get("name", ""),
                       color=PALETTE[si % len(PALETTE)], zorder=3)
    if chart.get("x_label"):
        ax.set_xlabel(chart["x_label"])
    if chart.get("y_label"):
        ax.set_ylabel(chart["y_label"])
    if series and series[0].get("name"):
        ax.legend()


def _funnel_chart(ax, chart):
    import numpy as np
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    values = [float(v) for v in (series[0].get("data", []) if series else [])]
    if not values:
        return
    y = np.arange(len(values))
    # 按值大小绘制宽度递减的横向条, 居中
    for yi, v in enumerate(values):
        ax.barh(yi, v, height=0.6, color=PALETTE[yi % len(PALETTE)], zorder=3)
        ax.text(v, yi, f" {v:g}", va="center", ha="left", fontsize=8)
    ax.set_yticks(y)
    ax.set_yticklabels(cats)
    ax.invert_yaxis()
    ax.set_xlim(0, max(values) * 1.25)


def _gauge_chart(ax, chart):
    import numpy as np
    series = chart.get("series") or []
    value = float(series[0].get("data", [0])[0]) if series and series[0].get("data") else 0
    vmin = float(chart.get("min", 0))
    vmax = float(chart.get("max", 100))
    theta = np.linspace(np.pi, 0, 100)
    r = 1.0
    ax.plot(np.cos(theta) * r, np.sin(theta) * r, color="#cbd5e1", linewidth=12, solid_capstyle="round")
    ratio = max(0.0, min(1.0, (value - vmin) / (vmax - vmin) if vmax > vmin else 0))
    n_on = int(ratio * 100)
    ax.plot(np.cos(theta[:max(n_on, 1)]) * r, np.sin(theta[:max(n_on, 1)]) * r,
            color="#2563eb", linewidth=12, solid_capstyle="round")
    ax.text(0, -0.25, f"{value:g}", ha="center", va="center", fontsize=28, fontweight="bold")
    ax.text(0, -0.55, chart.get("title") or "", ha="center", va="center", fontsize=10, color="#64748b")
    ax.set_xlim(-1.3, 1.3)
    ax.set_ylim(-0.8, 1.2)
    ax.axis("off")


_CHART_RENDERERS = {
    "bar": _bar_chart, "barh": _barh_chart, "line": _line_chart,
    "area": _area_chart, "pie": _pie_chart, "donut": lambda ax, c: _pie_chart(ax, c, donut=True),
    "radar": _radar_chart, "scatter": _scatter_chart, "funnel": _funnel_chart,
    "gauge": _gauge_chart,
}


def _chart_to_svg(chart):
    """渲染图表为 SVG 字符串 (矢量, 供 HTML)。失败返回 None。"""
    if not _mpl_setup():
        return None
    try:
        import matplotlib.pyplot as plt
        fig = plt.figure()
        ax = _chart_axes(chart, fig)
        renderer = _CHART_RENDERERS.get(chart.get("type", "bar"))
        if not renderer:
            return None
        if chart.get("type") in ("radar", "pie", "donut", "gauge"):
            if chart.get("type") == "radar":
                import numpy as np
                ax.remove()
                ax = fig.add_subplot(111, polar=True)
            else:
                pass  # pie/donut/gauge 用默认矩形 ax 但 axis equal/off
        renderer(ax, chart)
        if chart.get("title"):
            ax.set_title(chart["title"])
        fig.tight_layout()
        buf = io.StringIO()
        fig.savefig(buf, format="svg")
        plt.close(fig)
        svg = buf.getvalue()
        m = re.search(r"<svg[\s\S]*</svg>", svg)
        return m.group(0) if m else None
    except Exception:
        return None


def _chart_to_png(chart, dpi=220):
    """渲染图表为 PNG bytes (供 PDF/DOCX/XLSX/PPTX)。失败返回 None。"""
    if not _mpl_setup():
        return None
    try:
        import matplotlib.pyplot as plt
        fig = plt.figure()
        ax = _chart_axes(chart, fig)
        renderer = _CHART_RENDERERS.get(chart.get("type", "bar"))
        if not renderer:
            return None
        if chart.get("type") == "radar":
            import numpy as np
            ax.remove()
            ax = fig.add_subplot(111, polar=True)
        renderer(ax, chart)
        if chart.get("title"):
            ax.set_title(chart["title"])
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=dpi)
        plt.close(fig)
        return buf.getvalue()
    except Exception:
        return None


# ======================================================================
# HTML 渲染器 (自包含, 专业排版, 可打印 PDF)
# ======================================================================

def md_to_html(md_text, output_path):
    blocks = _parse_blocks(md_text)
    title = _extract_title(blocks)

    # 收集标题用于目录
    toc_entries = []
    for typ, payload in blocks:
        if typ == "heading" and payload[0] <= 2:
            toc_entries.append((payload[0], _strip_inline(payload[1])))
    has_toc = len(toc_entries) >= 3

    parts = []
    parts.append("<!DOCTYPE html>\n<html lang=\"zh-CN\">\n<head>\n<meta charset=\"utf-8\">")
    parts.append("<meta name=\"viewport\" content=\"width=device-width, initial-scale=1\">")
    parts.append(f"<title>{title}</title>")
    parts.append(f"<style>{_HTML_CSS}</style>\n</head>\n<body>")

    # 封面
    parts.append("<header class=\"cover\">")
    parts.append(f"<h1>{title}</h1>")
    parts.append("<p class=\"subtitle\">Meerkat AI · 智能报告</p>")
    parts.append("<span class=\"meta\">由 AI 自动生成 · 数据图表</span>")
    parts.append("</header>")

    # 目录
    if has_toc:
        parts.append("<nav class=\"toc\"><h2>目录</h2><ul>")
        for idx, (lvl, txt) in enumerate(toc_entries):
            cls = "lvl2" if lvl == 2 else ""
            parts.append(f"<li><a class=\"{cls}\" href=\"#sec-{idx}\">{txt}</a></li>")
        parts.append("</ul></nav>")

    # 正文
    parts.append("<main>")
    h_idx = 0
    for typ, payload in blocks:
        if typ == "heading":
            lvl, txt = payload
            if lvl <= 2:
                anchor = f"sec-{h_idx}"
                h_idx += 1
                parts.append(f"<h{min(lvl,6)} id=\"{anchor}\">{_inline_to_html(txt)}</h{min(lvl,6)}>")
            else:
                parts.append(f"<h{min(lvl,6)}>{_inline_to_html(txt)}</h{min(lvl,6)}>")
        elif typ == "para":
            parts.append(f"<p>{_inline_to_html(payload)}</p>")
        elif typ == "list":
            parts.append("<ul>")
            for item in payload:
                parts.append(f"<li>{_inline_to_html(item)}</li>")
            parts.append("</ul>")
        elif typ == "quote":
            parts.append(f"<blockquote>{_inline_to_html(payload)}</blockquote>")
        elif typ == "hr":
            parts.append("<hr>")
        elif typ == "code":
            parts.append(f"<pre><code>{_escape_html(payload)}</code></pre>")
        elif typ == "image":
            alt, fmt, b64 = payload
            parts.append(f"<img src=\"data:image/{fmt};base64,{b64}\" alt=\"{alt}\">")
        elif typ == "table":
            header, data = _parse_table_rows(payload)
            parts.append("<table>")
            if header:
                parts.append("<thead><tr>" + "".join(f"<th>{_inline_to_html(c)}</th>" for c in header) + "</tr></thead>")
            parts.append("<tbody>")
            for row in data:
                parts.append("<tr>" + "".join(f"<td>{_inline_to_html(c)}</td>" for c in row) + "</tr>")
            parts.append("</tbody></table>")
        elif typ == "chart":
            svg = _chart_to_svg(payload)
            if svg:
                parts.append(f"<div class=\"figure\">{svg}</div>")
            else:
                parts.append(f"<p><em>[图表: {payload.get('title', '')}]</em></p>")
    parts.append("</main>\n</body>\n</html>")

    html = "\n".join(parts)
    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)
    return output_path


def _escape_html(text):
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# ======================================================================
# PDF 渲染器 (reportlab, 封面/目录/页眉页脚 + Noto CJK + 高清图表)
# ======================================================================

def md_to_pdf(md_text, output_path):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                    TableStyle, PageBreak, Image)
    from reportlab.platypus.tableofcontents import TableOfContents
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfbase.pdfmetrics import registerFontFamily

    blocks = _parse_blocks(md_text)
    title = _extract_title(blocks)

    # ---- 注册中文字体 (TrueType 轮廓) ----
    reg, bold, sreg, sbold = _resolve_pdf_fonts()

    def _register(fname, path):
        try:
            if path and os.path.exists(path):
                if path.lower().endswith(".ttc"):
                    pdfmetrics.registerFont(TTFont(fname, path, subfontIndex=0))
                else:
                    pdfmetrics.registerFont(TTFont(fname, path))
                return fname
        except Exception:
            pass
        return None

    _sans = _register("CJKSans", reg) or "Helvetica"
    _sans_b = _register("CJKSansBold", bold) or _sans
    _serif = _register("CJKSerif", sreg) or _sans
    _serif_b = _register("CJKSerifBold", sbold) or _serif

    try:
        registerFontFamily("CJKSans", normal="CJKSans", bold="CJKSansBold",
                           italic="CJKSans", boldItalic="CJKSansBold")
    except Exception:
        pass
    try:
        registerFontFamily("CJKSerif", normal="CJKSerif", bold="CJKSerifBold",
                           italic="CJKSerif", boldItalic="CJKSerifBold")
    except Exception:
        pass

    # ---- 样式 ----
    styles = getSampleStyleSheet()
    body_font = "CJKSans" if "CJKSans" in pdfmetrics.getRegisteredFontNames() else "Helvetica"
    heading_font = body_font

    styles.add(ParagraphStyle("CNTitle", fontName=heading_font, fontSize=26, leading=34,
                              textColor=colors.white, alignment=1, spaceAfter=12))
    styles.add(ParagraphStyle("CNH1", fontName=heading_font, fontSize=18, leading=26,
                              textColor=colors.HexColor("#1e3a8a"), spaceBefore=18, spaceAfter=10))
    styles.add(ParagraphStyle("CNH2", fontName=heading_font, fontSize=15, leading=22,
                              textColor=colors.HexColor("#2563eb"), spaceBefore=14, spaceAfter=8))
    styles.add(ParagraphStyle("CNH3", fontName=heading_font, fontSize=13, leading=20,
                              textColor=colors.HexColor("#334155"), spaceBefore=10, spaceAfter=6))
    styles.add(ParagraphStyle("CNBody", fontName=body_font, fontSize=10.5, leading=18,
                              textColor=colors.HexColor("#334155"), spaceAfter=8))
    styles.add(ParagraphStyle("CNList", fontName=body_font, fontSize=10.5, leading=18,
                              leftIndent=16, bulletIndent=4, spaceAfter=4,
                              textColor=colors.HexColor("#334155")))
    styles.add(ParagraphStyle("CNQuote", fontName=body_font, fontSize=10, leading=16,
                              leftIndent=14, textColor=colors.HexColor("#64748b"),
                              borderPadding=8, backColor=colors.HexColor("#f8fafc"),
                              borderColor=colors.HexColor("#0ea5e9"), borderWidth=0, spaceAfter=10))

    # ---- 封面 / 页眉页脚回调 (局部函数, 闭包捕获 title/body_font) ----
    def _draw_cover(canvas, doc_):
        w, h = A4
        canvas.saveState()
        canvas.setFillColor(colors.HexColor("#1e3a8a"))
        canvas.rect(0, 0, w, h, fill=1, stroke=0)
        canvas.setFillColor(colors.white)
        bold_name = "CJKSansBold" if "CJKSansBold" in pdfmetrics.getRegisteredFontNames() else "Helvetica-Bold"
        canvas.setFont(bold_name, 30)
        canvas.drawCentredString(w / 2, h / 2 + 0.6 * cm, title)
        canvas.setFont(body_font, 13)
        canvas.drawCentredString(w / 2, h / 2 - 1.2 * cm, "Meerkat AI · 智能报告")
        canvas.restoreState()

    def _draw_header_footer(canvas, doc_):
        w, h = A4
        canvas.saveState()
        canvas.setFont(body_font, 8)
        canvas.setFillColor(colors.HexColor("#94a3b8"))
        canvas.drawString(2 * cm, h - 1.4 * cm, title)
        canvas.setStrokeColor(colors.HexColor("#e2e8f0"))
        canvas.line(2 * cm, h - 1.6 * cm, w - 2 * cm, h - 1.6 * cm)
        canvas.drawCentredString(w / 2, 1.0 * cm, f"— {doc_.page} —")
        canvas.restoreState()

    class DocTemplate(SimpleDocTemplate):
        def afterFlowable(self, flowable):
            if isinstance(flowable, Paragraph):
                st = flowable.style.name
                text = flowable.getPlainText()
                if st in ("CNH1", "CNH2", "CNH3"):
                    lvl = {"CNH1": 0, "CNH2": 1, "CNH3": 2}.get(st, 0)
                    key = f"h-{self.page}-{len(text)}"
                    self.canv.bookmarkPage(key)
                    self.canv.addOutlineEntry(text, key, level=lvl, closed=False)
                    self.notify("TOCEntry", (lvl, text, self.page, key))

    doc = DocTemplate(output_path, pagesize=A4,
                      leftMargin=2 * cm, rightMargin=2 * cm,
                      topMargin=2.2 * cm, bottomMargin=2.2 * cm,
                      title=title, author="Meerkat AI")

    story = []
    # 封面独占第一页
    story.append(PageBreak())

    # 目录 (需 multiBuild)
    toc = TableOfContents()
    toc.levelStyles = [
        ParagraphStyle(name="TOC1", fontName=heading_font, fontSize=11, leading=18, leftIndent=6),
        ParagraphStyle(name="TOC2", fontName=heading_font, fontSize=10, leading=16, leftIndent=24),
        ParagraphStyle(name="TOC3", fontName=heading_font, fontSize=9, leading=14, leftIndent=42),
    ]
    story.append(Paragraph("目录", styles["CNH1"]))
    story.append(toc)
    story.append(PageBreak())

    # 正文
    for typ, payload in blocks:
        if typ == "heading":
            lvl, txt = payload
            st = {1: "CNH1", 2: "CNH2", 3: "CNH3"}.get(lvl, "CNH3")
            story.append(Paragraph(_escape(txt), styles[st]))
        elif typ == "para":
            story.append(Paragraph(_escape(payload), styles["CNBody"]))
        elif typ == "list":
            for item in payload:
                story.append(Paragraph(_escape(item), styles["CNList"], bulletText="•"))
        elif typ == "quote":
            story.append(Paragraph(_escape(payload), styles["CNQuote"]))
        elif typ == "hr":
            story.append(Spacer(1, 8))
        elif typ == "code":
            story.append(Paragraph(_escape(payload), styles["CNQuote"]))
        elif typ == "image":
            _alt, _fmt, b64 = payload
            try:
                img = Image(io.BytesIO(base64.b64decode(b64)), width=14 * cm, height=10 * cm, kind="proportional")
                story.append(img)
            except Exception:
                story.append(Paragraph(_escape(_alt or "[图片]"), styles["CNBody"]))
        elif typ == "table":
            header, data = _parse_table_rows(payload)
            rows = ([header] if header else []) + data
            if rows:
                ncols = max(len(r) for r in rows)
                tdata = [[Paragraph(_escape(r[ci] if ci < len(r) else ""), styles["CNBody"]) for ci in range(ncols)] for r in rows]
                t = Table(tdata, colWidths=[(17 * cm) / ncols] * ncols)
                style = [
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cbd5e1")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ("TOPPADDING", (0, 0), (-1, -1), 5),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ]
                if header:
                    style.append(("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2563eb")))
                t.setStyle(TableStyle(style))
                story.append(t)
        elif typ == "chart":
            png = _chart_to_png(payload)
            if png:
                try:
                    img = Image(io.BytesIO(png), width=15 * cm, height=9 * cm, kind="proportional")
                    story.append(img)
                except Exception:
                    story.append(Paragraph(_escape(payload.get("title", "[图表]")), styles["CNBody"]))
        story.append(Spacer(1, 4))

    doc.multiBuild(story, onFirstPage=_draw_cover, onLaterPages=_draw_header_footer)
    return output_path


def _escape(text):
    """转义 reportlab 的 XML 特殊字符。"""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# ======================================================================
# DOCX 渲染器 (增强样式 + 图表 PNG 嵌入)
# ======================================================================

def md_to_docx(md_text, output_path):
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    # 标题颜色
    for typ, payload in _parse_blocks(md_text):
        if typ == "heading":
            level, text = payload
            h = doc.add_heading(_strip_inline(text), level=min(level, 4))
            for run in h.runs:
                run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
        elif typ == "para":
            doc.add_paragraph(_strip_inline(payload))
        elif typ == "image":
            from docx.shared import Inches
            _alt, _fmt, b64 = payload
            try:
                doc.add_picture(io.BytesIO(base64.b64decode(b64)), width=Inches(5.5))
            except Exception:
                doc.add_paragraph(_alt or "[图片]")
        elif typ == "list":
            for item in payload:
                doc.add_paragraph(_strip_inline(item), style="List Bullet")
        elif typ == "quote":
            p = doc.add_paragraph(_strip_inline(payload))
            p.paragraph_format.left_indent = Inches(0.3)
            for run in p.runs:
                run.font.italic = True
                run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        elif typ == "code":
            p = doc.add_paragraph(payload)
            for run in p.runs:
                run.font.name = "Courier New"
                run.font.size = Pt(9)
        elif typ == "table":
            header, data = _parse_table_rows(payload)
            rows = ([header] if header else []) + data
            if rows:
                table = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
                table.style = "Light Grid Accent 1"
                for ri, row in enumerate(rows):
                    for ci in range(len(table.columns)):
                        table.cell(ri, ci).text = row[ci] if ci < len(row) else ""
        elif typ == "chart":
            png = _chart_to_png(payload)
            if png:
                try:
                    doc.add_picture(io.BytesIO(png), width=Inches(6.0))
                except Exception:
                    doc.add_paragraph(payload.get("title", "[图表]"))
    doc.save(output_path)
    return output_path


# ======================================================================
# XLSX 渲染器 (多 Sheet + 样式 + 原生图表对象)
# ======================================================================

def md_to_xlsx(md_text, output_path):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference

    wb = Workbook()
    ws = wb.active
    ws.title = "报告"

    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill("solid", fgColor="2563EB")
    title_font = Font(bold=True, size=14, color="1E3A8A")
    h2_font = Font(bold=True, size=12, color="2563EB")
    h3_font = Font(bold=True, size=11, color="334155")
    cell_align = Alignment(vertical="top", wrap_text=True)
    header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    thin = Side(style="thin", color="CBD5E1")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    row_idx = 1
    chart_row = 0  # 记录图表数据起始行 (用于原生图表引用)

    def _write_table(header, data):
        nonlocal row_idx, chart_row
        rows = ([header] if header else []) + data
        if not rows:
            return
        start = row_idx
        for r in rows:
            for ci, cell in enumerate(r, start=1):
                c = ws.cell(row=row_idx, column=ci, value=cell)
                c.border = border
                c.alignment = cell_align
            row_idx += 1
        ncols = max(len(r) for r in rows)
        if header:
            for ci in range(1, ncols + 1):
                c = ws.cell(row=start, column=ci)
                c.font = header_font
                c.fill = header_fill
                c.alignment = header_align
            ws.freeze_panes = f"A{start + 1}"
        for ci in range(1, ncols + 1):
            max_len = max((len(str(r[ci - 1])) if ci - 1 < len(r) else 0) for r in rows)
            ws.column_dimensions[get_column_letter(ci)].width = min(max_len + 4, 40)
        chart_row = start
        row_idx += 1

    for typ, payload in _parse_blocks(md_text):
        if typ == "heading":
            level, text = payload
            c = ws.cell(row=row_idx, column=1, value=_strip_inline(text))
            c.font = title_font if level == 1 else (h2_font if level == 2 else h3_font)
            row_idx += 1
        elif typ == "table":
            header, data = _parse_table_rows(payload)
            _write_table(header, data)
        elif typ == "para":
            ws.cell(row=row_idx, column=1, value=_strip_inline(payload)).alignment = cell_align
            row_idx += 1
        elif typ == "list":
            for item in payload:
                ws.cell(row=row_idx, column=1, value="• " + _strip_inline(item)).alignment = cell_align
                row_idx += 1
        elif typ == "chart":
            # 把图表数据写入新 sheet, 并生成原生图表对象
            chart = payload
            cats = chart.get("categories") or []
            series = chart.get("series") or []
            if cats and series:
                cs = wb.create_sheet(title=(chart.get("title") or "图表")[:28])
                cs.cell(row=1, column=1, value="类别")
                for ci, cat in enumerate(cats, start=2):
                    cs.cell(row=1, column=ci, value=cat)
                for si, s in enumerate(series):
                    cs.cell(row=2 + si, column=1, value=s.get("name", f"系列{si+1}"))
                    for di, v in enumerate(s.get("data", []), start=2):
                        cs.cell(row=2 + si, column=di, value=float(v))
                ctype = chart.get("type", "bar")
                if ctype in ("pie", "donut"):
                    xc = PieChart()
                    data_ref = Reference(cs, min_col=2, min_row=2, max_row=1 + len(series))
                    cats_ref = Reference(cs, min_col=2, max_col=1 + len(cats), min_row=1, max_row=1)
                    xc.add_data(data_ref, titles_from_data=False)
                    xc.set_categories(cats_ref)
                elif ctype == "line":
                    xc = LineChart()
                    data_ref = Reference(cs, min_col=1, min_row=1, max_col=1 + len(cats), max_row=1 + len(series))
                    xc.add_data(data_ref, titles_from_data=True)
                    xc.set_categories(Reference(cs, min_col=2, min_row=1, max_col=1 + len(cats), max_row=1))
                else:
                    xc = BarChart()
                    xc.type = "col"
                    data_ref = Reference(cs, min_col=1, min_row=1, max_col=1 + len(cats), max_row=1 + len(series))
                    xc.add_data(data_ref, titles_from_data=True)
                    xc.set_categories(Reference(cs, min_col=2, min_row=1, max_col=1 + len(cats), max_row=1))
                xc.title = chart.get("title") or ""
                xc.height = 9
                xc.width = 16
                cs.add_chart(xc, f"A{3 + len(series)}")
                # 样式
                for ci in range(1, 1 + len(cats) + 1):
                    c = cs.cell(row=1, column=ci)
                    c.font = header_font
                    c.fill = header_fill
                    c.alignment = header_align
                    cs.column_dimensions[get_column_letter(ci)].width = 14
        elif typ == "image":
            from openpyxl.drawing.image import Image as XLImage
            _alt, _fmt, b64 = payload
            try:
                img = XLImage(io.BytesIO(base64.b64decode(b64)))
                img.width = 360
                img.height = 360
                ws.add_image(img, f"A{row_idx}")
                row_idx += 22
            except Exception:
                ws.cell(row=row_idx, column=1, value=_alt or "[图片]")
                row_idx += 1

    wb.save(output_path)
    return output_path


# ======================================================================
# PPTX 渲染器 (版式 + 原生图表对象)
# ======================================================================

def md_to_pptx(md_text, output_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blocks = _parse_blocks(md_text)
    title = _extract_title(blocks)

    def _add_title_bar(slide, text, sub=None):
        # 顶部色条 + 标题
        from pptx.enum.shapes import MSO_SHAPE
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
        bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8))
        tf = tb.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(1.6))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    sub_tb = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11.3), Inches(0.6))
    sub_tb.text_frame.text = "Meerkat AI · 智能报告"
    sub_tb.text_frame.paragraphs[0].font.size = Pt(18)
    sub_tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    sub_tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    pending_heading = None
    for typ, payload in blocks:
        if typ == "heading":
            lvl, txt = payload
            if lvl == 1:
                # 章节页
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _add_title_bar(slide, _strip_inline(txt))
                continue
            # 内容页标题
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_title_bar(slide, _strip_inline(txt))
            pending_heading = _strip_inline(txt)
        elif typ == "chart":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_title_bar(slide, chart_title := (payload.get("title") or "数据图表"))
            cats = payload.get("categories") or []
            series = payload.get("series") or []
            if cats and series:
                cd = CategoryChartData()
                cd.categories = cats
                for s in series:
                    cd.add_series(s.get("name", "系列"), [float(v) for v in s.get("data", [])])
                ctype = payload.get("type", "bar")
                chart_type = {
                    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                    "barh": XL_CHART_TYPE.BAR_CLUSTERED,
                    "line": XL_CHART_TYPE.LINE_MARKERS,
                    "area": XL_CHART_TYPE.AREA,
                    "pie": XL_CHART_TYPE.PIE,
                    "donut": XL_CHART_TYPE.DOUGHNUT,
                }.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)
                gf = slide.shapes.add_chart(chart_type, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), cd)
                gf.chart.has_title = False
                gf.chart.has_legend = True
        elif typ == "table":
            header, data = _parse_table_rows(payload)
            rows = ([header] if header else []) + data
            if rows:
                ncols = max(len(r) for r in rows)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _add_title_bar(slide, pending_heading or "数据表")
                gf = slide.shapes.add_table(len(rows), ncols, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.5 * len(rows) + 0.4))
                tbl = gf.table
                for ri, row in enumerate(rows):
                    for ci in range(ncols):
                        tbl.cell(ri, ci).text = row[ci] if ci < len(row) else ""
                        tbl.cell(ri, ci).text_frame.paragraphs[0].font.size = Pt(12)
        elif typ == "list":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_title_bar(slide, pending_heading or "要点")
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, item in enumerate(payload):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = "▪ " + _strip_inline(item)
                p.font.size = Pt(18)
                p.space_after = Pt(10)
        elif typ == "para" and len(_strip_inline(payload)) > 0:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_title_bar(slide, pending_heading or "内容")
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.text = _strip_inline(payload)
            tf.paragraphs[0].font.size = Pt(18)
        elif typ == "image":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _alt, _fmt, b64 = payload
            try:
                slide.shapes.add_picture(io.BytesIO(base64.b64decode(b64)), Inches(1), Inches(1.3), width=Inches(11.3))
            except Exception:
                pass

    if len(prs.slides._sldIdLst) == 0:
        prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(output_path)
    return output_path


# ======================================================================
# 统一入口
# ======================================================================

_FORMATS = {
    "html": md_to_html,
    "docx": md_to_docx,
    "pdf": md_to_pdf,
    "xlsx": md_to_xlsx,
    "pptx": md_to_pptx,
}


def convert(fmt, md_text, output_path):
    """统一转换入口: convert('html'|'docx'|'pdf'|'xlsx'|'pptx', markdown_text, output_path)。"""
    fmt = fmt.lower().lstrip(".")
    if fmt not in _FORMATS:
        raise ValueError(f"不支持的格式: {fmt}, 支持: {list(_FORMATS)}")
    return _FORMATS[fmt](md_text, output_path)


if __name__ == "__main__":
    if len(sys.argv) != 4:
        print("用法: python3 doc_tools.py <md2html|md2docx|md2pdf|md2xlsx|md2pptx> <input.md> <output>")
        sys.exit(1)
    cmd, inp, outp = sys.argv[1], sys.argv[2], sys.argv[3]
    fmt = cmd.replace("md2", "")
    with open(inp, encoding="utf-8") as f:
        md = f.read()
    result = convert(fmt, md, outp)
    print(f"[done] {fmt} → {result}")
