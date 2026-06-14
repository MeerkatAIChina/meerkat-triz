"""
TRIZ-raw 语料库构建器
将PDF/DOCX/PPTX/DOC原始材料提取、分块、metadata化，输出JSONL语料。
"""

import json
import logging
import os
import time
import hashlib
from abc import ABC, abstractmethod
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Dict, List, Optional, Set, Tuple, Any, Iterator

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@dataclass
class Page:
    """单个页面/幻灯片/章节提取结果"""
    text: str
    page_num: int = 1
    heading: Optional[str] = None


@dataclass
class ExtractedDocument:
    """单个文件提取结果"""
    source_path: str
    title: Optional[str] = None
    category: Optional[str] = None
    file_type: Optional[str] = None
    pages: List[Page] = field(default_factory=list)


class Extractor(ABC):
    """文件提取器基类"""

    @abstractmethod
    def extract(self, path: Path) -> ExtractedDocument:
        """从文件中提取文本，返回ExtractedDocument"""
        pass


class PDFExtractor(Extractor):
    """PDF提取器：优先提取可选中文字，可选OCR回退"""

    def __init__(self, ocr_enabled: bool = True, min_text_chars: int = 20):
        self.ocr_enabled = ocr_enabled
        self.min_text_chars = min_text_chars

    def extract(self, path: Path) -> ExtractedDocument:
        import fitz  # pymupdf

        doc = ExtractedDocument(
            source_path=str(path),
            file_type="pdf",
        )

        try:
            pdf = fitz.open(str(path))
        except Exception as e:
            logger.error(f"无法打开PDF {path}: {e}")
            return doc

        for page_idx in range(len(pdf)):
            page = pdf.load_page(page_idx)
            text = page.get_text().strip()

            # OCR回退：页面文字太少时尝试OCR
            if self.ocr_enabled and len(text) < self.min_text_chars:
                try:
                    from PIL import Image
                    import pytesseract
                    pix = page.get_pixmap(dpi=200)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text = pytesseract.image_to_string(img, lang="chi_sim+eng").strip()
                except Exception as e:
                    logger.warning(f"OCR失败 {path} 第{page_idx + 1}页: {e}")

            if text:
                doc.pages.append(Page(text=text, page_num=page_idx + 1))

        pdf.close()
        return doc


class DocxExtractor(Extractor):
    """DOCX提取器：按段落和标题提取"""

    def extract(self, path: Path) -> ExtractedDocument:
        from docx import Document

        doc = ExtractedDocument(source_path=str(path), file_type="docx")

        try:
            document = Document(str(path))
        except Exception as e:
            logger.error(f"无法打开DOCX {path}: {e}")
            return doc

        doc.title = document.core_properties.title or None
        current_heading = None

        for para in document.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # 识别标题样式
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                current_heading = text
                continue

            doc.pages.append(Page(text=text, page_num=1, heading=current_heading))

        return doc


class PptxExtractor(Extractor):
    """PPTX提取器：每页幻灯片作为一个Page"""

    def extract(self, path: Path) -> ExtractedDocument:
        from pptx import Presentation

        doc = ExtractedDocument(source_path=str(path), file_type="pptx")

        try:
            prs = Presentation(str(path))
        except Exception as e:
            logger.error(f"无法打开PPTX {path}: {e}")
            return doc

        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            text = "\n".join(texts)
            if text:
                doc.pages.append(Page(text=text, page_num=slide_idx))

        return doc


@dataclass
class Chunk:
    """分块后的语料记录"""
    text: str
    source_path: str
    category: Optional[str]
    file_type: Optional[str]
    page_num: Optional[int]
    heading: Optional[str]
    chunk_index: int
    token_count: int


class SemanticChunker:
    """语义分块器：按标题/段落边界合并短段落到目标token数"""

    def __init__(
        self,
        target_tokens: int = 2048,
        max_tokens: int = 4096,
        chars_per_token: float = 1.0,
        min_chars: int = 50,
    ):
        self.target_tokens = target_tokens
        self.max_tokens = max_tokens
        self.chars_per_token = chars_per_token
        self.min_chars = min_chars

    def _estimate_tokens(self, text: str) -> int:
        return max(1, int(len(text) / self.chars_per_token))

    def _split_into_segments(self, text: str) -> List[str]:
        """按空行分割成语义段"""
        segments = []
        current = []
        for line in text.splitlines():
            stripped = line.strip()
            if not stripped:
                if current:
                    segments.append("\n".join(current))
                    current = []
            else:
                current.append(stripped)
        if current:
            segments.append("\n".join(current))
        return [s for s in segments if len(s) >= self.min_chars]

    def chunk_document(self, doc: ExtractedDocument, chunk_index_start: int = 0) -> Tuple[List[Chunk], int]:
        """对一个ExtractedDocument进行分块，返回Chunk列表和下一个chunk_index"""
        chunks = []
        chunk_index = chunk_index_start
        current_texts: List[str] = []
        current_tokens = 0
        current_heading: Optional[str] = None
        current_page: Optional[int] = None

        def flush():
            nonlocal chunks, chunk_index, current_texts, current_tokens, current_heading, current_page
            if not current_texts:
                return
            text = "\n\n".join(current_texts)
            if len(text) >= self.min_chars:
                chunks.append(Chunk(
                    text=text,
                    source_path=doc.source_path,
                    category=doc.category,
                    file_type=doc.file_type,
                    page_num=current_page,
                    heading=current_heading,
                    chunk_index=chunk_index,
                    token_count=self._estimate_tokens(text),
                ))
                chunk_index += 1
            current_texts = []
            current_tokens = 0

        for page in doc.pages:
            segments = self._split_into_segments(page.text)
            if not segments:
                continue

            # 记录当前页的起始heading
            if page.heading:
                current_heading = page.heading
            if page.page_num:
                current_page = page.page_num

            for seg in segments:
                seg_tokens = self._estimate_tokens(seg)

                # 单个段超过max_tokens则强行截断
                if seg_tokens > self.max_tokens:
                    flush()
                    chunks.append(Chunk(
                        text=seg[: int(self.max_tokens * self.chars_per_token)],
                        source_path=doc.source_path,
                        category=doc.category,
                        file_type=doc.file_type,
                        page_num=current_page,
                        heading=current_heading,
                        chunk_index=chunk_index,
                        token_count=self.max_tokens,
                    ))
                    chunk_index += 1
                    continue

                # 合并到当前chunk会超过目标token数，先flush
                if current_tokens + seg_tokens > self.target_tokens and current_texts:
                    flush()

                current_texts.append(seg)
                current_tokens += seg_tokens

                # 如果当前heading来自这个segment的上一段标题，更新heading
                if page.heading and not current_heading:
                    current_heading = page.heading

        flush()
        return chunks, chunk_index


class CorpusWriter:
    """语料库写入器：输出JSONL、统计信息和失败文件列表"""

    def __init__(
        self,
        output_dir: str,
        output_filename: str = "triz_corpus.jsonl",
        stats_filename: str = "triz_corpus_stats.json",
        failed_files_filename: str = "failed_files.json",
        deduplicate: bool = True,
    ):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.output_path = self.output_dir / output_filename
        self.stats_path = self.output_dir / stats_filename
        self.failed_files_path = self.output_dir / failed_files_filename
        self.deduplicate = deduplicate

    def _deduplicate(self, chunks: List[Chunk]) -> List[Chunk]:
        if not self.deduplicate:
            return chunks
        seen: Set[str] = set()
        unique = []
        for c in chunks:
            key = hashlib.md5(c.text.encode("utf-8")).hexdigest()
            if key not in seen:
                seen.add(key)
                unique.append(c)
        removed = len(chunks) - len(unique)
        if removed > 0:
            logger.info(f"去重: 移除 {removed} 个重复chunk")
        return unique

    def write(
        self,
        chunks: List[Chunk],
        failed_files: Optional[List[Dict[str, str]]] = None,
    ) -> Dict[str, Any]:
        chunks = self._deduplicate(chunks)

        # 写JSONL
        with open(self.output_path, "w", encoding="utf-8") as f:
            for idx, chunk in enumerate(chunks):
                record = {
                    "id": f"triz-raw-{idx:08d}",
                    "text": chunk.text,
                    "metadata": {
                        "source_path": chunk.source_path,
                        "category": chunk.category,
                        "file_type": chunk.file_type,
                        "page_num": chunk.page_num,
                        "heading": chunk.heading,
                        "chunk_index": chunk.chunk_index,
                        "token_count": chunk.token_count,
                    },
                }
                f.write(json.dumps(record, ensure_ascii=False) + "\n")

        # 统计信息
        category_counts: Dict[Optional[str], int] = {}
        file_type_counts: Dict[Optional[str], int] = {}
        token_counts = [c.token_count for c in chunks]

        for c in chunks:
            category_counts[c.category] = category_counts.get(c.category, 0) + 1
            file_type_counts[c.file_type] = file_type_counts.get(c.file_type, 0) + 1

        stats = {
            "total_records": len(chunks),
            "total_tokens": sum(token_counts),
            "avg_tokens": round(sum(token_counts) / len(token_counts), 2) if token_counts else 0,
            "max_tokens": max(token_counts) if token_counts else 0,
            "min_tokens": min(token_counts) if token_counts else 0,
            "category_distribution": category_counts,
            "file_type_distribution": file_type_counts,
        }

        with open(self.stats_path, "w", encoding="utf-8") as f:
            json.dump(stats, f, ensure_ascii=False, indent=2)

        # 失败文件（始终写入，避免残留旧失败列表）
        failed_files = failed_files or []
        with open(self.failed_files_path, "w", encoding="utf-8") as f:
            json.dump(failed_files, f, ensure_ascii=False, indent=2)

        logger.info(f"语料库写入完成: {len(chunks)} 条记录 -> {self.output_path}")
        return stats


def build_corpus(
    raw_dir: str,
    output_dir: str,
    output_filename: str = "triz_corpus.jsonl",
    stats_filename: str = "triz_corpus_stats.json",
    failed_files_filename: str = "failed_files.json",
    supported_extensions: Optional[List[str]] = None,
    skip_extensions: Optional[List[str]] = None,
    chunk_target_tokens: int = 2048,
    chunk_max_tokens: int = 4096,
    chars_per_token: float = 1.0,
    min_chars: int = 50,
    deduplicate: bool = True,
    ocr_enabled: bool = True,
    ocr_min_text_chars: int = 20,
    resume: bool = True,
) -> Dict[str, Any]:
    """
    构建TRIZ语料库主入口

    Args:
        raw_dir: TRIZ-raw目录路径
        output_dir: 输出目录
        output_filename: 输出JSONL文件名
        stats_filename: 统计文件名
        failed_files_filename: 失败文件列表名
        supported_extensions: 支持提取的扩展名
        skip_extensions: 默认跳过的扩展名
        chunk_target_tokens: 目标chunk token数
        chunk_max_tokens: 最大chunk token数
        chars_per_token: 字符到token的估算比例
        min_chars: 最小chunk字符数
        deduplicate: 是否去重
        ocr_enabled: 是否启用OCR
        ocr_min_text_chars: OCR触发阈值
        resume: 是否跳过已处理的文件

    Returns:
        统计信息字典
    """
    raw_dir = Path(raw_dir)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    supported_extensions = supported_extensions or [".pdf", ".docx", ".pptx", ".doc"]
    skip_extensions = set(skip_extensions or [".jpg", ".jpeg", ".png", ".gif", ".bmp",
                                                ".webp", ".jfif", ".mov", ".mp4", ".avi",
                                                ".mkv", ".zip", ".rar", ".7z", ".xlsx",
                                                ".xls", ".csv"])

    # 已处理文件集合 (用于断点续跑)
    processed_sources: Set[str] = set()
    output_path = output_dir / output_filename
    if resume and output_path.exists():
        with open(output_path, "r", encoding="utf-8") as f:
            for line in f:
                try:
                    record = json.loads(line)
                    processed_sources.add(record["metadata"]["source_path"])
                except Exception:
                    pass
        logger.info(f"断点续跑: 已跳过 {len(processed_sources)} 个已处理文件")

    # 提取器映射
    extractors = {
        ".pdf": PDFExtractor(ocr_enabled=ocr_enabled, min_text_chars=ocr_min_text_chars),
        ".docx": DocxExtractor(),
        ".pptx": PptxExtractor(),
        ".doc": DocxExtractor(),  # .doc用python-docx尝试，失败则跳过
    }

    chunker = SemanticChunker(
        target_tokens=chunk_target_tokens,
        max_tokens=chunk_max_tokens,
        chars_per_token=chars_per_token,
        min_chars=min_chars,
    )
    writer = CorpusWriter(
        output_dir=str(output_dir),
        output_filename=output_filename,
        stats_filename=stats_filename,
        failed_files_filename=failed_files_filename,
        deduplicate=deduplicate,
    )

    all_chunks: List[Chunk] = []
    failed_files: List[Dict[str, str]] = []
    chunk_index_counter = 0

    if not raw_dir.exists():
        raise FileNotFoundError(f"原始目录不存在: {raw_dir}")

    # 收集所有文件
    files = []
    for ext in supported_extensions:
        files.extend(raw_dir.rglob(f"*{ext}"))

    # 过滤跳过扩展名
    files = [f for f in files if f.suffix.lower() not in skip_extensions]

    logger.info(f"发现 {len(files)} 个待处理文件")

    for file_path in files:
        rel_path = str(file_path.relative_to(raw_dir))
        source_path = f"TRIZ-raw/{rel_path}"

        if source_path in processed_sources:
            logger.info(f"跳过已处理: {source_path}")
            continue

        ext = file_path.suffix.lower()
        extractor = extractors.get(ext)
        if extractor is None:
            logger.warning(f"无可用提取器: {file_path}")
            continue

        try:
            doc = extractor.extract(file_path)
            doc.source_path = source_path
            doc.category = file_path.parent.relative_to(raw_dir).parts[0] if file_path.parent != raw_dir else "root"

            chunks, chunk_index_counter = chunker.chunk_document(doc, chunk_index_start=chunk_index_counter)
            all_chunks.extend(chunks)
            logger.info(f"处理完成: {source_path} -> {len(chunks)} chunks")
        except Exception as e:
            logger.error(f"处理失败 {source_path}: {e}")
            failed_files.append({"path": source_path, "error": str(e)})

    stats = writer.write(all_chunks, failed_files=failed_files)
    return stats
